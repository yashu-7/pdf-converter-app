import os
import shutil
import zipfile
import uuid  # For creating unique session folders
from io import BytesIO
from flask import Flask, render_template, request, jsonify, send_from_directory, url_for
from werkzeug.utils import secure_filename # For security
from dotenv import load_dotenv

# --- All Conversion Libraries ---
import pdf2docx
import fitz  # PyMuPDF
from pypdf import PdfWriter, PdfReader
from pptx import Presentation
from pptx.util import Inches

# This line loads your .env file for local testing. It does nothing on Render.
load_dotenv()

# --- Environment Variable & Secret Setup ---
GA_ID_KEY = 'GA_MEASUREMENT_ID'
FORMSPREE_ID_KEY = 'FORMSPREE_FORM_ID'

GA_ID = os.environ.get(GA_ID_KEY)
FORMSPREE_ID = os.environ.get(FORMSPREE_ID_KEY)
FORMSPREE_URL = f"https://formspree.io/f/{FORMSPREE_ID}" if FORMSPREE_ID else None

# --- Application Setup ---
app = Flask(__name__)
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMP_FOLDER = os.path.join(BASE_DIR, 'temp_files')

# A dictionary to hold details about our tools
TOOLS = {
    "pdf-merger": {"title": "PDF Merger", "description": "Combine multiple PDFs into one document.", "icon": "🔗"},
    "split-pdf": {"title": "Split PDF", "description": "Split a single PDF into separate pages or at a specific page.", "icon": "✂️"},
    "pdf-to-word": {"title": "PDF to Word", "description": "Convert a PDF to an editable .docx file.", "icon": "📄"},
    "pdf-to-ppt": {"title": "PDF to PowerPoint", "description": "Convert a PDF into a .pptx presentation.", "icon": "🖼️"}
}

# ==============================================================================
# SECTION 1: CORE CONVERSION FUNCTIONS (As provided in your reference)
# ==============================================================================

def merge_pdfs(input_paths, output_path):
    merger = PdfWriter()
    for path in input_paths:
        merger.append(path)
    with open(output_path, "wb") as f_out:
        merger.write(f_out)

def split_pdf_all_pages(input_path, output_folder):
    reader = PdfReader(input_path)
    os.makedirs(output_folder, exist_ok=True)
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        with open(os.path.join(output_folder, f"page_{i + 1}.pdf"), "wb") as f_out:
            writer.write(f_out)
    return len(reader.pages)

def split_pdf_at_page(input_path, output_folder, split_page):
    reader = PdfReader(input_path)
    total_pages = len(reader.pages)
    if not (1 <= split_page < total_pages):
        raise ValueError(f"Invalid split page. Must be between 1 and {total_pages - 1}.")
    os.makedirs(output_folder, exist_ok=True)
    # Part 1
    writer1 = PdfWriter()
    for i in range(split_page):
        writer1.add_page(reader.pages[i])
    with open(os.path.join(output_folder, f"Part1_Pages_1-{split_page}.pdf"), "wb") as f_out:
        writer1.write(f_out)
    # Part 2
    writer2 = PdfWriter()
    for i in range(split_page, total_pages):
        writer2.add_page(reader.pages[i])
    with open(os.path.join(output_folder, f"Part2_Pages_{split_page + 1}-{total_pages}.pdf"), "wb") as f_out:
        writer2.write(f_out)

def convert_pdf_to_word(input_path, output_path):
    cv = pdf2docx.Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()

def convert_pdf_to_ppt(input_path, output_path):
    pdf_doc = fitz.open(input_path)
    prs = Presentation()
    for page in pdf_doc:
        pix = page.get_pixmap(dpi=150)
        image_bytes = pix.tobytes("png")
        # Set slide dimensions based on page aspect ratio
        prs.slide_width = Inches(10) # Standard 16:9 width
        aspect_ratio = page.rect.width / page.rect.height if page.rect.height > 0 else 1.77
        prs.slide_height = int(prs.slide_width / aspect_ratio)
        slide_layout = prs.slide_layouts[6] # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(BytesIO(image_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)


# ==============================================================================
# SECTION 2: FLASK ROUTES
# ==============================================================================

@app.context_processor
def inject_global_variables():
    """Makes variables available to all templates automatically."""
    return dict(ga_id=GA_ID, formspree_url=FORMSPREE_URL)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/<tool_id>')
def tool_page(tool_id):
    if tool_id not in TOOLS:
        return "Tool not found!", 404
    return render_template('tool_page.html', tool=TOOLS[tool_id], tool_id=tool_id)


@app.route('/upload', methods=['POST'])
def upload_files():
    # Create a unique session folder for this request to prevent file conflicts
    session_id = str(uuid.uuid4())
    session_folder = os.path.join(TEMP_FOLDER, session_id)
    os.makedirs(session_folder, exist_ok=True)
    
    tool_id = request.form.get('tool_id')
    files = request.files.getlist('files[]')

    if not files or files[0].filename == '':
        return jsonify({'status': 'error', 'message': 'No files were selected.'}), 400

    try:
        # Save uploaded files securely
        saved_paths = []
        for f in files:
            # Sanitize filename for security
            filename = secure_filename(f.filename)
            path = os.path.join(session_folder, filename)
            f.save(path)
            saved_paths.append(path)

        # --- Tool-specific logic ---
        if tool_id == 'pdf-merger':
            if len(saved_paths) < 2:
                raise ValueError('Please select at least two files to merge.')
            output_filename = "merged_document.pdf"
            output_path = os.path.join(session_folder, output_filename)
            merge_pdfs(saved_paths, output_path)
            message = f"Successfully merged {len(files)} files."

        elif tool_id == 'split-pdf':
            split_output_folder = os.path.join(session_folder, "split_results")
            split_option = request.form.get('split-option')
            output_filename = "split_pages.zip"

            if split_option == 'range':
                split_page = request.form.get('split_page', type=int)
                if not split_page: raise ValueError('Split page not provided for range split.')
                split_pdf_at_page(saved_paths[0], split_output_folder, split_page)
                message = f"PDF successfully split at page {split_page}."
            else: # all-pages
                num_pages = split_pdf_all_pages(saved_paths[0], split_output_folder)
                message = f"Successfully split the PDF into {num_pages} individual pages."
            
            # Zip the split results
            zip_path = os.path.join(session_folder, output_filename)
            with zipfile.ZipFile(zip_path, 'w') as zf:
                for root, _, zfiles in os.walk(split_output_folder):
                    for zfile in zfiles:
                        zf.write(os.path.join(root, zfile), arcname=zfile)

        elif tool_id == 'pdf-to-word':
            base_filename = os.path.splitext(os.path.basename(saved_paths[0]))[0]
            output_filename = f"{base_filename}.docx"
            output_path = os.path.join(session_folder, output_filename)
            convert_pdf_to_word(saved_paths[0], output_path)
            message = "Successfully converted PDF to Word."

        elif tool_id == 'pdf-to-ppt':
            base_filename = os.path.splitext(os.path.basename(saved_paths[0]))[0]
            output_filename = f"{base_filename}.pptx"
            output_path = os.path.join(session_folder, output_filename)
            convert_pdf_to_ppt(saved_paths[0], output_path)
            message = "Successfully converted PDF to PowerPoint."
        else:
            raise ValueError("Invalid tool specified.")

        download_url = url_for('download_file', session_id=session_id, filename=output_filename)
        return jsonify({'status': 'success', 'message': message, 'download_url': download_url})

    except Exception as e:
        # Cleanup the session folder on error
        shutil.rmtree(session_folder, ignore_errors=True)
        app.logger.error(f"Error in {tool_id}: {e}", exc_info=True)
        return jsonify({'status': 'error', 'message': f'An error occurred: {e}'}), 500


@app.route('/download/<session_id>/<filename>')
def download_file(session_id, filename):
    """Securely sends the file for download from the correct session folder."""
    directory = os.path.join(TEMP_FOLDER, session_id)
    try:
        # After sending, the session folder could be scheduled for cleanup, but for now we leave it.
        # On Render, these files are temporary anyway.
        return send_from_directory(directory, filename, as_attachment=True)
    except FileNotFoundError:
        return "File not found. It may have been deleted.", 404


if __name__ == '__main__':
    # Ensure the base temp folder exists
    if not os.path.exists(TEMP_FOLDER):
        os.makedirs(TEMP_FOLDER)
    app.run(debug=False)
